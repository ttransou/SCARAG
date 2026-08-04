from __future__ import annotations

import json
import re
import zipfile
from importlib import import_module
from email import policy
from email.parser import BytesParser
from datetime import UTC, datetime
from pathlib import Path
from typing import Any

from bs4 import BeautifulSoup
from docx import Document as DocxDocument
from pptx import Presentation as PptxPresentation
from pypdf import PdfReader
from openpyxl import load_workbook
from openpyxl.utils.exceptions import InvalidFileException

try:
    import pdfplumber
except ImportError:  # pragma: no cover - optional at runtime, required in package deps
    pdfplumber = None  # type: ignore[assignment]

try:
    xlrd = import_module("xlrd")
except ImportError:  # pragma: no cover - optional at runtime, required in package deps
    xlrd = None


def _utc_now_iso() -> str:
    return datetime.now(UTC).replace(microsecond=0).isoformat().replace("+00:00", "Z")


def _read_text_file(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="ignore")


def _flatten_json(value: Any, prefix: str = "") -> list[str]:
    lines: list[str] = []
    if isinstance(value, dict):
        if not value and prefix:
            lines.append(f"{prefix}: {{}}")
            return lines
        for key, item in value.items():
            child_prefix = f"{prefix}.{key}" if prefix else str(key)
            lines.extend(_flatten_json(item, child_prefix))
        return lines

    if isinstance(value, list):
        if not value:
            lines.append(f"{prefix}: []" if prefix else "[]")
            return lines
        for index, item in enumerate(value):
            child_prefix = f"{prefix}[{index}]" if prefix else f"[{index}]"
            lines.extend(_flatten_json(item, child_prefix))
        return lines

    scalar = json.dumps(value, ensure_ascii=False) if value is not None else "null"
    if prefix:
        lines.append(f"{prefix}: {scalar}")
    else:
        lines.append(str(scalar))
    return lines


def _parse_html(path: Path) -> str:
    html = _read_text_file(path)
    soup = BeautifulSoup(html, "html.parser")
    for tag in soup(["script", "style"]):
        tag.decompose()
    return "\n\n".join(part.strip() for part in re.split(r"\n{2,}", soup.get_text("\n", strip=True)) if part.strip())


def _parse_xml(path: Path) -> str:
    raw = _read_text_file(path)
    soup = BeautifulSoup(raw, "xml")
    return "\n\n".join(part.strip() for part in re.split(r"\n{2,}", soup.get_text("\n", strip=True)) if part.strip())


def _parse_legacy_doc(path: Path) -> str:
    raw = path.read_bytes()
    decoded = raw.decode("latin-1", errors="ignore").replace("\x00", " ")
    candidates = re.findall(r"[A-Za-z0-9][A-Za-z0-9\s,.;:()/_\-]{3,}", decoded)
    cleaned_lines = [re.sub(r"\s+", " ", part).strip() for part in candidates]
    cleaned_lines = [part for part in cleaned_lines if part]
    if cleaned_lines:
        return "\n".join(cleaned_lines)

    fallback = re.sub(r"[^\x20-\x7E\n\t]", " ", decoded)
    return re.sub(r"\s+", " ", fallback).strip()


def _parse_docx(path: Path) -> tuple[str, list[dict[str, Any]], list[dict[str, Any]]]:
    document = DocxDocument(path)

    paragraphs = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]
    table_blocks: list[str] = []
    table_metadata: list[dict[str, Any]] = []
    image_markers: list[dict[str, Any]] = []

    for table_index, table in enumerate(document.tables):
        rows: list[list[str]] = []
        for row in table.rows:
            cleaned = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cleaned:
                rows.append(cleaned)

        if not rows:
            continue

        header = rows[0]
        body = rows[1:]
        lines = [" | ".join(header)] + [" | ".join(row) for row in body]
        table_blocks.append("\n".join(lines))
        table_metadata.append(
            {
                "table_id": f"docx_table_{table_index}",
                "row_count": len(rows),
                "column_count": len(header),
                "header_fields": header,
            }
        )

    # Detect inline and anchored images by scanning relationship targets.
    for rel_id, relationship in document.part.rels.items():
        target_ref = getattr(relationship, "target_ref", "") or ""
        if "image" not in target_ref.lower():
            continue
        image_markers.append(
            {
                "marker_id": f"docx_image_{len(image_markers)}",
                "location": f"relationship:{rel_id}",
                "content_type": "image",
                "target": target_ref,
            }
        )

    combined = "\n\n".join([*paragraphs, *table_blocks]).strip()
    return combined, table_metadata, image_markers


def _parse_pdf(path: Path) -> tuple[str, list[dict[str, Any]], list[dict[str, Any]]]:
    text_parts: list[str] = []
    table_metadata: list[dict[str, Any]] = []
    image_markers: list[dict[str, Any]] = []

    if pdfplumber is not None:
        with pdfplumber.open(str(path)) as pdf:
            for page_index, page in enumerate(pdf.pages):
                page_text = (page.extract_text() or "").strip()
                if page_text:
                    text_parts.append(page_text)

                tables = page.extract_tables() or []
                for table_index, table_rows in enumerate(tables):
                    cleaned_rows = [
                        [str(cell).strip() for cell in row if cell is not None and str(cell).strip()]
                        for row in table_rows
                        if row
                    ]
                    cleaned_rows = [row for row in cleaned_rows if row]
                    if not cleaned_rows:
                        continue

                    header = cleaned_rows[0]
                    body = cleaned_rows[1:]
                    text_parts.append("\n".join([" | ".join(header), *[" | ".join(row) for row in body]]))
                    table_metadata.append(
                        {
                            "table_id": f"pdf_page_{page_index}_table_{table_index}",
                            "row_count": len(cleaned_rows),
                            "column_count": len(header),
                            "header_fields": header,
                            "page_index": page_index,
                        }
                    )

                for image_index, image in enumerate(page.images or []):
                    image_markers.append(
                        {
                            "marker_id": f"pdf_page_{page_index}_image_{image_index}",
                            "location": f"page:{page_index}",
                            "content_type": "image",
                            "bbox": [image.get("x0"), image.get("top"), image.get("x1"), image.get("bottom")],
                        }
                    )

    if text_parts:
        return "\n\n".join(part for part in text_parts if part.strip()), table_metadata, image_markers

    reader = PdfReader(str(path))
    pages = [page.extract_text() or "" for page in reader.pages]
    joined = "\n\n".join(page.strip() for page in pages if page.strip())
    return joined, table_metadata, image_markers


def _parse_pptx(path: Path) -> tuple[str, list[dict[str, Any]], list[dict[str, Any]]]:
    presentation = PptxPresentation(path)
    slide_blocks: list[str] = []
    table_metadata: list[dict[str, Any]] = []
    image_markers: list[dict[str, Any]] = []
    for slide_index, slide in enumerate(presentation.slides):
        texts = [shape.text.strip() for shape in slide.shapes if getattr(shape, "text", "").strip()]

        for shape_index, shape in enumerate(slide.shapes):
            if not getattr(shape, "has_table", False):
                continue
            table = shape.table
            rows: list[list[str]] = []
            for row in table.rows:
                cleaned = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cleaned:
                    rows.append(cleaned)

            if not rows:
                continue

            header = rows[0]
            body = rows[1:]
            lines = [" | ".join(header)] + [" | ".join(row) for row in body]
            texts.append("\n".join(lines))
            table_metadata.append(
                {
                    "table_id": f"pptx_slide_{slide_index}_table_{shape_index}",
                    "row_count": len(rows),
                    "column_count": len(header),
                    "header_fields": header,
                    "slide_index": slide_index,
                }
            )

            if hasattr(shape, "image"):
                image_markers.append(
                    {
                        "marker_id": f"pptx_slide_{slide_index}_shape_{shape_index}",
                        "location": f"slide:{slide_index}",
                        "content_type": "image",
                    }
                )

        for shape_index, shape in enumerate(slide.shapes):
            if not hasattr(shape, "image"):
                continue
            image_markers.append(
                {
                    "marker_id": f"pptx_slide_{slide_index}_image_{shape_index}",
                    "location": f"slide:{slide_index}",
                    "content_type": "image",
                }
            )

        if texts:
            slide_blocks.append("\n".join(texts))
    return "\n\n".join(slide_blocks), table_metadata, image_markers


def _looks_like_spreadsheet_header(cells: list[str]) -> bool:
    if len(cells) < 2:
        return False

    alpha_cells = sum(1 for cell in cells if any(char.isalpha() for char in cell))
    numeric_like_cells = sum(
        1
        for cell in cells
        if cell.replace(".", "", 1).isdigit() or cell.replace("-", "", 1).isdigit()
    )
    return alpha_cells >= max(1, len(cells) // 2) and numeric_like_cells == 0


def _parse_xlsx(path: Path) -> tuple[str, list[dict[str, Any]]]:
    workbook = load_workbook(path, read_only=True, data_only=True)
    sheets = []
    table_metadata: list[dict[str, Any]] = []
    line_cursor = 1
    for sheet in workbook.worksheets:
        rows = []
        row_cells: list[list[str]] = []
        for row in sheet.iter_rows(values_only=True):
            cleaned = [str(value).strip() for value in row if value is not None and str(value).strip()]
            if cleaned:
                row_cells.append(cleaned)
                rows.append(" | ".join(cleaned))
        if rows:
            sheets.append("\n".join(rows))

            has_header = _looks_like_spreadsheet_header(row_cells[0])
            header_fields = row_cells[0] if has_header else []
            line_start_index = line_cursor
            line_end_index = line_cursor + len(rows) - 1
            table_metadata.append(
                {
                    "table_id": f"xlsx_{sheet.title}",
                    "row_count": len(rows),
                    "column_count": len(row_cells[0]),
                    "header_fields": header_fields,
                    "has_header": has_header,
                    "sheet_name": sheet.title,
                    "line_start_index": line_start_index,
                    "line_end_index": line_end_index,
                    "data_row_start_index": 2 if has_header else 1,
                    "data_row_end_index": len(rows),
                }
            )
            line_cursor = line_end_index + 1
    return "\n\n".join(sheets), table_metadata


def _parse_xls_with_xlrd(path: Path) -> tuple[str, list[dict[str, Any]]]:
    if xlrd is None:
        raise RuntimeError("xlrd is required to parse .xls files")

    workbook = xlrd.open_workbook(str(path))
    sheets = []
    table_metadata: list[dict[str, Any]] = []
    line_cursor = 1
    for sheet in workbook.sheets():
        rows = []
        row_cells: list[list[str]] = []
        for row_index in range(sheet.nrows):
            cleaned = [
                str(value).strip()
                for value in sheet.row_values(row_index)
                if str(value).strip()
            ]
            if cleaned:
                row_cells.append(cleaned)
                rows.append(" | ".join(cleaned))
        if rows:
            sheets.append("\n".join(rows))

            has_header = _looks_like_spreadsheet_header(row_cells[0])
            header_fields = row_cells[0] if has_header else []
            line_start_index = line_cursor
            line_end_index = line_cursor + len(rows) - 1
            table_metadata.append(
                {
                    "table_id": f"xls_{sheet.name}",
                    "row_count": len(rows),
                    "column_count": len(row_cells[0]),
                    "header_fields": header_fields,
                    "has_header": has_header,
                    "sheet_name": sheet.name,
                    "line_start_index": line_start_index,
                    "line_end_index": line_end_index,
                    "data_row_start_index": 2 if has_header else 1,
                    "data_row_end_index": len(rows),
                }
            )
            line_cursor = line_end_index + 1
    return "\n\n".join(sheets), table_metadata


def _parse_spreadsheet(path: Path) -> tuple[str, list[dict[str, Any]]]:
    suffix = path.suffix.lower()
    if suffix == ".xls":
        try:
            return _parse_xlsx(path)
        except InvalidFileException:
            return _parse_xls_with_xlrd(path)
    return _parse_xlsx(path)


def _parse_result(payload: Any) -> tuple[str, list[dict[str, Any]], list[dict[str, Any]]]:
    if isinstance(payload, tuple):
        if len(payload) == 3:
            text = str(payload[0])
            table_metadata = payload[1] if isinstance(payload[1], list) else []
            image_markers = payload[2] if isinstance(payload[2], list) else []
            return text, table_metadata, image_markers
        if len(payload) == 2:
            text = str(payload[0])
            table_metadata = payload[1] if isinstance(payload[1], list) else []
            return text, table_metadata, []
    return str(payload), [], []


def _parse_json(path: Path) -> str:
    data = json.loads(_read_text_file(path))
    return "\n".join(_flatten_json(data))


def _parse_csv(path: Path) -> str:
    return _read_text_file(path)


def _parse_mhtml(path: Path) -> str:
    raw = path.read_bytes()
    try:
        message = BytesParser(policy=policy.default).parsebytes(raw)
    except Exception:
        return _read_text_file(path)

    html_parts: list[str] = []
    plain_parts: list[str] = []

    for part in message.walk():
        if part.is_multipart():
            continue
        content_type = part.get_content_type().lower()
        payload = part.get_payload(decode=True) or b""
        charset = part.get_content_charset() or "utf-8"
        decoded = payload.decode(charset, errors="ignore").strip()
        if not decoded:
            continue

        if content_type == "text/html":
            soup = BeautifulSoup(decoded, "html.parser")
            html_parts.append(soup.get_text("\n", strip=True))
        elif content_type == "text/plain":
            plain_parts.append(decoded)

    if html_parts:
        return "\n\n".join(part for part in html_parts if part.strip())
    if plain_parts:
        return "\n\n".join(part for part in plain_parts if part.strip())
    return _read_text_file(path)


def _empty_ingestion_diagnostics() -> dict[str, Any]:
    return {
        "files": [],
        "summary": {
            "total_files": 0,
            "loaded_files": 0,
            "skipped_files": 0,
            "unsupported_files": 0,
            "internal_skipped_files": 0,
            "parser_counts": {},
            "skip_reason_counts": {},
        },
    }


def _record_ingestion_file(
    diagnostics: dict[str, Any],
    *,
    source: str,
    parser: str,
    status: str,
    skip_reason: str | None = None,
) -> None:
    files = diagnostics.get("files")
    summary = diagnostics.get("summary")
    if not isinstance(files, list) or not isinstance(summary, dict):
        return

    record = {
        "source": source,
        "parser": parser,
        "status": status,
        "skip_reason": skip_reason,
    }
    files.append(record)

    summary["total_files"] = int(summary.get("total_files", 0)) + 1
    parser_counts = summary.setdefault("parser_counts", {})
    if isinstance(parser_counts, dict):
        parser_counts[parser] = int(parser_counts.get(parser, 0)) + 1

    if status == "loaded":
        summary["loaded_files"] = int(summary.get("loaded_files", 0)) + 1
        return

    summary["skipped_files"] = int(summary.get("skipped_files", 0)) + 1
    if skip_reason == "unsupported_format":
        summary["unsupported_files"] = int(summary.get("unsupported_files", 0)) + 1
    if skip_reason == "internal_artifact":
        summary["internal_skipped_files"] = int(summary.get("internal_skipped_files", 0)) + 1

    if skip_reason:
        skip_reason_counts = summary.setdefault("skip_reason_counts", {})
        if isinstance(skip_reason_counts, dict):
            skip_reason_counts[skip_reason] = int(skip_reason_counts.get(skip_reason, 0)) + 1


def _is_internal_ingestion_artifact(path: Path) -> bool:
    name = path.name
    return name in {
        ".scarag_lifecycle_state.json",
        ".scarag_lifecycle_audit.jsonl",
    }


def _load_documents_internal(data_path: str | Path) -> tuple[list[dict[str, Any]], dict[str, Any]]:
    root = Path(data_path)
    diagnostics = _empty_ingestion_diagnostics()
    if not root.exists():
        return [], diagnostics

    documents: list[dict[str, Any]] = []
    parser_by_suffix: dict[str, tuple[Any, str]] = {
        ".txt": (_read_text_file, "text_file_parser"),
        ".md": (_read_text_file, "text_file_parser"),
        ".json": (_parse_json, "json_parser"),
        ".csv": (_parse_csv, "csv_parser"),
        ".mhtml": (_parse_mhtml, "mhtml_parser"),
        ".mht": (_parse_mhtml, "mhtml_parser"),
        ".html": (_parse_html, "html_parser"),
        ".htm": (_parse_html, "html_parser"),
        ".xml": (_parse_xml, "xml_parser"),
        ".doc": (_parse_legacy_doc, "doc_legacy_parser"),
    }

    for path in sorted(root.rglob("*")):
        if not path.is_file():
            continue

        if _is_internal_ingestion_artifact(path):
            _record_ingestion_file(
                diagnostics,
                source=str(path),
                parser="internal_artifact",
                status="skipped",
                skip_reason="internal_artifact",
            )
            continue

        suffix = path.suffix.lower()
        if suffix in parser_by_suffix:
            parser, extraction_method = parser_by_suffix[suffix]
            try:
                text = parser(path)
            except Exception:
                _record_ingestion_file(
                    diagnostics,
                    source=str(path),
                    parser=extraction_method,
                    status="skipped",
                    skip_reason="parse_error",
                )
                continue
            documents.append(
                {
                    "source": str(path),
                    "text": text,
                    "doc_type": "unknown",
                    "extraction_method": extraction_method,
                    "extraction_ts": _utc_now_iso(),
                }
            )
            _record_ingestion_file(
                diagnostics,
                source=str(path),
                parser=extraction_method,
                status="loaded",
            )
            continue

        if suffix == ".docx":
            text, table_metadata, image_markers = _parse_result(_parse_docx(path))
            documents.append(
                {
                    "source": str(path),
                    "text": text,
                    "doc_type": "unknown",
                    "extraction_method": "docx_parser",
                    "extraction_ts": _utc_now_iso(),
                    "table_metadata": table_metadata,
                    "image_markers": image_markers,
                }
            )
            _record_ingestion_file(
                diagnostics,
                source=str(path),
                parser="docx_parser",
                status="loaded",
            )
            continue

        if suffix == ".pdf":
            text, table_metadata, image_markers = _parse_result(_parse_pdf(path))
            documents.append(
                {
                    "source": str(path),
                    "text": text,
                    "doc_type": "unknown",
                    "extraction_method": "pdf_parser",
                    "extraction_ts": _utc_now_iso(),
                    "table_metadata": table_metadata,
                    "image_markers": image_markers,
                }
            )
            _record_ingestion_file(
                diagnostics,
                source=str(path),
                parser="pdf_parser",
                status="loaded",
            )
            continue

        if suffix == ".pptx":
            text, table_metadata, image_markers = _parse_result(_parse_pptx(path))
            documents.append(
                {
                    "source": str(path),
                    "text": text,
                    "doc_type": "unknown",
                    "extraction_method": "pptx_parser",
                    "extraction_ts": _utc_now_iso(),
                    "table_metadata": table_metadata,
                    "image_markers": image_markers,
                }
            )
            _record_ingestion_file(
                diagnostics,
                source=str(path),
                parser="pptx_parser",
                status="loaded",
            )
            continue

        if suffix == ".xlsx" or suffix == ".xls":
            text, table_metadata, image_markers = _parse_result(_parse_spreadsheet(path))
            extraction_method = "xls_parser" if suffix == ".xls" else "xlsx_parser"
            documents.append(
                {
                    "source": str(path),
                    "text": text,
                    "doc_type": "unknown",
                    "extraction_method": extraction_method,
                    "extraction_ts": _utc_now_iso(),
                    "table_metadata": table_metadata,
                    "image_markers": image_markers,
                }
            )
            _record_ingestion_file(
                diagnostics,
                source=str(path),
                parser=extraction_method,
                status="loaded",
            )
            continue

        _record_ingestion_file(
            diagnostics,
            source=str(path),
            parser="unsupported_parser",
            status="skipped",
            skip_reason="unsupported_format",
        )

    return documents, diagnostics


def load_documents_with_diagnostics(data_path: str | Path) -> tuple[list[dict[str, Any]], dict[str, Any]]:
    return _load_documents_internal(data_path)


def load_documents(data_path: str | Path) -> list[dict[str, Any]]:
    """Load a simple list of document-like records from a folder.

    The repository README describes a richer ingestion pipeline, but this
    reference implementation keeps the interface lightweight until real parsers
    are added.
    """
    documents, _diagnostics = _load_documents_internal(data_path)
    return documents
