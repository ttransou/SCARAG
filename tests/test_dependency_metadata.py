import zipfile
import base64
from pathlib import Path
from tempfile import TemporaryDirectory
from types import SimpleNamespace

from docx import Document as DocxDocument
from pptx import Presentation as PptxPresentation
from pptx.util import Inches
from openpyxl import Workbook
from openpyxl.utils.exceptions import InvalidFileException

import tomllib

import scarag.ingestion.loader as loader_module
from scarag.ingestion.loader import load_documents


REPO_ROOT = Path(__file__).resolve().parents[1]
VALID_PNG_BASE64 = (
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMCAO2r8R8AAAAASUVORK5CYII="
)


def test_document_parser_dependencies_are_declared() -> None:
    requirements = (REPO_ROOT / "requirements.txt").read_text(encoding="utf-8")
    pyproject = tomllib.loads((REPO_ROOT / "pyproject.toml").read_text(encoding="utf-8"))

    declared_requirements = pyproject["project"]["dependencies"]

    for package in [
        "pypdf",
        "pdfplumber",
        "python-docx",
        "python-pptx",
        "openpyxl",
        "xlrd",
        "beautifulsoup4",
        "lxml",
    ]:
        assert package in requirements
        assert any(package in dependency for dependency in declared_requirements)


def test_loader_handles_json_csv_and_mhtml_documents() -> None:
    with TemporaryDirectory() as temp_dir:
        temp_path = Path(temp_dir)
        (temp_path / "sample.json").write_text(
            '{"title": "Example", "value": 42}', encoding="utf-8"
        )
        (temp_path / "sample.csv").write_text(
            "name,value\nalpha,1\n", encoding="utf-8"
        )
        (temp_path / "sample.mhtml").write_bytes(
            b"From: sender@example.com\n"
            b"Content-Type: multipart/alternative; boundary=boundary\n\n"
            b"--boundary\n"
            b"Content-Type: text/plain; charset=utf-8\n\n"
            b"Plain text body\n"
            b"--boundary\n"
            b"Content-Type: text/html; charset=utf-8\n\n"
            b"<html><body><p>MHTML body</p></body></html>\n"
            b"--boundary--\n"
        )

        documents = load_documents(temp_path)

    assert len(documents) == 3
    assert any(doc["source"].endswith("sample.json") for doc in documents)
    assert any(doc["source"].endswith("sample.csv") for doc in documents)
    assert any(doc["source"].endswith("sample.mhtml") for doc in documents)
    assert any("MHTML body" in doc["text"] for doc in documents)


def test_loader_flattens_nested_json_paths() -> None:
    with TemporaryDirectory() as temp_dir:
        temp_path = Path(temp_dir)
        (temp_path / "nested.json").write_text(
            '{"user": {"name": "Alice", "roles": ["admin", "analyst"]}, "active": true}',
            encoding="utf-8",
        )

        documents = load_documents(temp_path)

    assert len(documents) == 1
    text = documents[0]["text"]
    assert "user.name: \"Alice\"" in text
    assert "user.roles[0]: \"admin\"" in text
    assert "user.roles[1]: \"analyst\"" in text
    assert "active: true" in text


def test_loader_parses_html_and_docx_content() -> None:
    with TemporaryDirectory() as temp_dir:
        temp_path = Path(temp_dir)
        (temp_path / "sample.html").write_text(
            "<html><body><h1>Heading</h1><p>HTML paragraph</p></body></html>",
            encoding="utf-8",
        )

        docx_path = temp_path / "sample.docx"
        with zipfile.ZipFile(docx_path, "w") as archive:
            archive.writestr(
                "[Content_Types].xml",
                """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
                <Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
                  <Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
                  <Default Extension="xml" ContentType="application/xml"/>
                  <Override PartName="/word/document.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>
                </Types>""",
            )
            archive.writestr(
                "_rels/.rels",
                """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
                <Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
                  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="word/document.xml"/>
                </Relationships>""",
            )
            archive.writestr(
                "word/document.xml",
                """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
                <w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">
                  <w:body>
                    <w:p><w:r><w:t>Docx paragraph</w:t></w:r></w:p>
                  </w:body>
                </w:document>""",
            )

        documents = load_documents(temp_path)

    assert len(documents) == 2
    html_doc = next(doc for doc in documents if doc["source"].endswith("sample.html"))
    docx_doc = next(doc for doc in documents if doc["source"].endswith("sample.docx"))
    assert "HTML paragraph" in html_doc["text"]
    assert "Docx paragraph" in docx_doc["text"]


def test_loader_parses_pptx_and_xlsx_content() -> None:
    with TemporaryDirectory() as temp_dir:
        temp_path = Path(temp_dir)

        pptx_path = temp_path / "sample.pptx"
        presentation = PptxPresentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = "Slide title"
        slide.placeholders[1].text = "Slide body"
        presentation.save(pptx_path)

        workbook = Workbook()
        sheet = workbook.active
        sheet.append(["Name", "Value"])
        sheet.append(["Alpha", 1])
        workbook.save(temp_path / "sample.xlsx")

        documents = load_documents(temp_path)

    assert len(documents) == 2
    pptx_doc = next(doc for doc in documents if doc["source"].endswith("sample.pptx"))
    xlsx_doc = next(doc for doc in documents if doc["source"].endswith("sample.xlsx"))
    assert "Slide title" in pptx_doc["text"]
    assert "Slide body" in pptx_doc["text"]
    assert "Name" in xlsx_doc["text"]
    assert "Alpha" in xlsx_doc["text"]
    assert isinstance(pptx_doc.get("table_metadata"), list)
    assert isinstance(xlsx_doc.get("table_metadata"), list)
    assert xlsx_doc["table_metadata"][0]["sheet_name"] == sheet.title
    assert xlsx_doc["table_metadata"][0]["line_start_index"] == 1
    assert xlsx_doc["table_metadata"][0]["line_end_index"] == 2
    assert xlsx_doc["table_metadata"][0]["has_header"] is True


def test_loader_uses_xls_fallback_parser_when_openpyxl_rejects_file(monkeypatch) -> None:
    with TemporaryDirectory() as temp_dir:
        temp_path = Path(temp_dir)
        xls_path = temp_path / "legacy.xls"
        xls_path.write_bytes(b"placeholder")

        def _raise_invalid(*args, **kwargs):
            raise InvalidFileException("legacy xls")

        monkeypatch.setattr(loader_module, "load_workbook", _raise_invalid)
        monkeypatch.setattr(
            loader_module,
            "_parse_xls_with_xlrd",
            lambda _: (
                "Name | Value\nAlpha | 1",
                [
                    {
                        "table_id": "xls_Sheet1",
                        "row_count": 2,
                        "column_count": 2,
                        "header_fields": ["Name", "Value"],
                        "has_header": True,
                        "sheet_name": "Sheet1",
                        "line_start_index": 1,
                        "line_end_index": 2,
                        "data_row_start_index": 2,
                        "data_row_end_index": 2,
                    }
                ],
            ),
        )

        documents = load_documents(temp_path)

    assert len(documents) == 1
    assert documents[0]["source"].endswith("legacy.xls")
    assert documents[0]["extraction_method"] == "xls_parser"
    assert "Alpha" in documents[0]["text"]
    assert documents[0]["table_metadata"][0]["sheet_name"] == "Sheet1"


def test_loader_extracts_docx_table_metadata() -> None:
    with TemporaryDirectory() as temp_dir:
        temp_path = Path(temp_dir)
        docx_path = temp_path / "table.docx"

        document = DocxDocument()
        table = document.add_table(rows=2, cols=2)
        table.cell(0, 0).text = "Name"
        table.cell(0, 1).text = "Value"
        table.cell(1, 0).text = "Alpha"
        table.cell(1, 1).text = "1"
        document.save(docx_path)

        documents = load_documents(temp_path)

    assert len(documents) == 1
    doc = documents[0]
    assert doc["source"].endswith("table.docx")
    assert "Name | Value" in doc["text"]
    assert "Alpha | 1" in doc["text"]
    assert doc["table_metadata"]
    assert doc["table_metadata"][0]["header_fields"] == ["Name", "Value"]
    assert isinstance(doc.get("image_markers"), list)


def test_loader_extracts_docx_image_markers() -> None:
    with TemporaryDirectory() as temp_dir:
        temp_path = Path(temp_dir)
        docx_path = temp_path / "image.docx"

        document = DocxDocument()
        document.add_paragraph("Document with image")
        image_path = temp_path / "sample.png"
        image_path.write_bytes(base64.b64decode(VALID_PNG_BASE64))
        document.add_picture(str(image_path))
        document.save(docx_path)

        documents = load_documents(temp_path)

    assert len(documents) == 1
    doc = next(item for item in documents if item["source"].endswith("image.docx"))
    assert doc["image_markers"]
    assert doc["image_markers"][0]["content_type"] == "image"


def test_loader_extracts_pptx_table_metadata() -> None:
    with TemporaryDirectory() as temp_dir:
        temp_path = Path(temp_dir)
        pptx_path = temp_path / "table.pptx"

        presentation = PptxPresentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        table_shape = slide.shapes.add_table(rows=2, cols=2, left=0, top=0, width=9144000, height=914400)
        table = table_shape.table
        table.cell(0, 0).text = "Name"
        table.cell(0, 1).text = "Value"
        table.cell(1, 0).text = "Beta"
        table.cell(1, 1).text = "2"
        presentation.save(pptx_path)

        documents = load_documents(temp_path)

    assert len(documents) == 1
    doc = documents[0]
    assert doc["source"].endswith("table.pptx")
    assert "Name | Value" in doc["text"]
    assert "Beta | 2" in doc["text"]
    assert doc["table_metadata"]
    assert doc["table_metadata"][0]["header_fields"] == ["Name", "Value"]
    assert isinstance(doc.get("image_markers"), list)


def test_loader_extracts_pptx_image_markers() -> None:
    with TemporaryDirectory() as temp_dir:
        temp_path = Path(temp_dir)
        pptx_path = temp_path / "image.pptx"
        image_path = temp_path / "sample.png"
        image_path.write_bytes(base64.b64decode(VALID_PNG_BASE64))

        presentation = PptxPresentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.add_picture(str(image_path), Inches(1), Inches(1), Inches(1), Inches(1))
        presentation.save(pptx_path)

        documents = load_documents(temp_path)

    assert len(documents) == 1
    doc = next(item for item in documents if item["source"].endswith("image.pptx"))
    assert doc["image_markers"]
    assert doc["image_markers"][0]["content_type"] == "image"


def test_loader_handles_nested_mhtml_multipart() -> None:
    with TemporaryDirectory() as temp_dir:
        temp_path = Path(temp_dir)
        nested = (
            b"From: sender@example.com\n"
            b"MIME-Version: 1.0\n"
            b"Content-Type: multipart/mixed; boundary=outer\n\n"
            b"--outer\n"
            b"Content-Type: multipart/alternative; boundary=inner\n\n"
            b"--inner\n"
            b"Content-Type: text/plain; charset=utf-8\n\n"
            b"Plain fallback\n"
            b"--inner\n"
            b"Content-Type: text/html; charset=utf-8\n\n"
            b"<html><body><p>Nested HTML body</p></body></html>\n"
            b"--inner--\n"
            b"--outer--\n"
        )
        (temp_path / "nested.mhtml").write_bytes(nested)

        documents = load_documents(temp_path)

    assert len(documents) == 1
    assert "Nested HTML body" in documents[0]["text"]


def test_loader_extracts_pdf_table_and_image_metadata_with_pdfplumber(monkeypatch) -> None:
    with TemporaryDirectory() as temp_dir:
        temp_path = Path(temp_dir)
        pdf_path = temp_path / "sample.pdf"
        pdf_path.write_bytes(b"%PDF-1.4\n%placeholder\n")

        fake_page = SimpleNamespace(
            extract_text=lambda: "Page text",
            extract_tables=lambda: [[("Name", "Value"), ("Alpha", "1")]],
            images=[{"x0": 1, "top": 2, "x1": 3, "bottom": 4}],
        )

        class _FakePdf:
            def __init__(self):
                self.pages = [fake_page]

            def __enter__(self):
                return self

            def __exit__(self, exc_type, exc, tb):
                return None

        monkeypatch.setattr(loader_module, "pdfplumber", SimpleNamespace(open=lambda *_: _FakePdf()))

        documents = load_documents(temp_path)

    assert len(documents) == 1
    doc = documents[0]
    assert "Page text" in doc["text"]
    assert "Name | Value" in doc["text"]
    assert doc["table_metadata"]
    assert doc["table_metadata"][0]["header_fields"] == ["Name", "Value"]
    assert doc["image_markers"]
    assert doc["image_markers"][0]["content_type"] == "image"
